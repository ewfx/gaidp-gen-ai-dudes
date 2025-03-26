#install : pip3 install python-pptx

from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Customer Dataset Validation"
subtitle.text = "Using OpenAI for Python Code Generation in Regulatory Validation"

# Slide 2: Dataset Overview
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Customer Dataset Overview"
content.text = (
    "The dataset used contains the following columns:\n"
    "- Customer_ID\n"
    "- Account_Balance\n"
    "- Transaction_Amount\n"
    "- Reported_Amount\n"
    "- Currency\n"
    "- Country\n"
    "- Transaction_Date\n"
    "- Risk_Score"
)

# Slide 3: Validation Workflow
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Validation Workflow"
content.text = (
    "1. Load the customer dataset.\n"
    "2. Input natural language validation rules.\n"
    "3. Use OpenAI's GPT to generate Python code dynamically.\n"
    "4. Execute the generated Python code to filter and validate the dataset.\n"
    "5. Handle and resolve errors during validation.\n"
    "6. Review filtered results for regulatory compliance."
)

# Slide 4: OpenAI-Powered Validation
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Role of OpenAI in Validation"
content.text = (
    "OpenAI simplifies validation by:\n"
    "- Accepting natural language instructions.\n"
    "- Dynamically generating Python code for dataset filtering.\n"
    "- Automating complex validation workflows.\n\n"
    "Example Rule:\n"
    "'Filter rows where Risk_Score > 3 and Account_Balance > 0.'"
)

# Slide 5: Code Implementation
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Python Code Example"
content.text = (
    "import openai\n"
    "openai.api_key = 'your-api-key'\n\n"
    "instructions = 'Filter rows where Risk_Score > 3 and Account_Balance > 0'\n\n"
    "prompt = f\"\"\"\n"
    "Based on the dataset, generate Python code to {instructions}.\n"
    "\"\"\"\n\n"
    "response = openai.ChatCompletion.create(\n"
    "    model='gpt-3.5-turbo',\n"
    "    messages=[\n"
    "        {'role': 'system', 'content': 'You are a Python assistant.'},\n"
    "        {'role': 'user', 'content': prompt}\n"
    "    ]\n"
    ")\n"
    "print(response['choices'][0]['message']['content'])"
)

# Slide 6: Example Dataset Preview
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Example Dataset (Preview)"
content.text = (
    "Customer_ID, Account_Balance, Transaction_Amount, Reported_Amount, Currency, Country, Transaction_Date, Risk_Score\n"
    "1001, 15000, 500, 500, USD, US, 2025-02-25, 3\n"
    "1002, 32000, 1200, 1200, EUR, DE, 2025-02-20, 2\n"
    "1003, -5000, 300, 300, GBP, UK, 2025-02-18, 6\n"
    "1004, 70000, 2000, 1800, USD, US, 2025-02-28, 5"
)

# Slide 7: Error Handling and Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Error Handling and Conclusion"
content.text = (
    "Error Handling:\n"
    "- OpenAI can suggest corrections for invalid instructions.\n"
    "- Dynamic execution of code allows immediate debugging.\n\n"
    "Conclusion:\n"
    "- OpenAI facilitates efficient and compliant validation.\n"
    "- Automation saves time and minimizes human error."
)

# Save the presentation
presentation.save("Customer_Data_Validation_PPT.pptx")

print("PowerPoint presentation has been created as 'Customer_Data_Validation_PPT.pptx'.")